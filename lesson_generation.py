import os
import re
import datetime
import socket
from pathlib import Path
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import traceback
import requests
import gspread
from oauth2client.service_account import ServiceAccountCredentials
from utils import read_file_contents, clean_text, extract_important_lines
from google import genai
import PyPDF2
import io

load_dotenv()

# Constants

# Class to store per-user session data
class UserSession:
    def __init__(self):
        self.user_id = None
        self.fileContent = None
        self.topic = None
        self.learning_outcomes = None
        self.user_directory = None
        self.sections = None

# Dictionary to store user sessions
user_sessions = {}

def get_user_session(user_id):
    """Get or create a user session for the given user_id"""
    if user_id not in user_sessions:
        user_sessions[user_id] = UserSession()
    return user_sessions[user_id]

role = None
main_prompt = None
questions_prompt = None
task_prompt = None
homework_prompt = None
revision_prompt = None
sections = None
topic_directory = None  # Path to the topic-specific directory
user_directory = None   # Path to the user-specific directory


DESIGN = '''
Design an HTML page using the provided data, ensuring a clean, structured layout with a modern and visually appealing design.  
- Use a vibrant and harmonious color palette to make the page visually engaging and easy to read.  
- Position the topic prominently at the top-center of the page.  
- Display the current date and day visibly on the page at all times.  
- Organize the content into distinct sections, with each section separated for clarity.  
- for each section use bollted list for the data that is related to that section.
- Add engaging and smooth animations throughout the page to elevate the user experience.  
- Ensure the page is styled correctly using CSS and HTML tags.
- Include CSS in the <head> section of  HTML document using the <style> tag.
- remove "Information 1:" or "Information 2:" and so on from the sections titles.
Focus on creating a professional and user-friendly interface.
'''

def add_row(version, time, topic, ip, computer_name):
    """Add a row to the Google Sheets tracking document."""
    scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
    creds = ServiceAccountCredentials.from_json_keyfile_name("gen-lang-client-0422190211-1c28abc28438.json", scope)
    client = gspread.authorize(creds)
    sheet = client.open("Lesson Generator users tracker").sheet1
    data = [version, time, topic, ip, computer_name]
    sheet.append_row(data)

##############
import requests

def ai_agent(prompt):
    try:
        print("Calling Mistral AI agent...")

        api_key = os.getenv("MISTRAL_API_KEY")
        if not api_key:
            raise ValueError("MISTRAL_API_KEY environment variable is not set")
        url = "https://api.mistral.ai/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "mistral-tiny",
            "messages": [{"role": "user", "content": prompt}]
        }

        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()

        return response.json()["choices"][0]["message"]["content"]

    except Exception as e:
        print(f"Error calling Mistral API: {e}")
        return None
###############



def ai_agent_google(prompt):
    print("Calling AI agent...")
    # Use Gemini API via HTTP requests
    print(prompt)

    gemini_key = os.getenv("GEMINI_API_KEY")
    if not gemini_key:
        raise ValueError("GEMINI_API_KEY environment variable is not set")
    os.environ['GEMINI_API_KEY'] = gemini_key

    

    client = genai.Client()
    response = client.models.generate_content(
    model="gemini-2.0-flash", 
    contents=prompt
    
  )
    
    return response.text
    print("AI agent call completed.")


def generate_topic_from_outcomes(outcomes):
    """Generate a topic based on learning outcomes."""
    prompt = f"""Based on these learning outcomes, suggest an appropriate topic that would cover all these learning objectives. Give me just the topic name without any explanation or additional text:

{outcomes}"""
    topic_name = ai_agent(prompt).strip().strip('"').strip("'").strip()
    
    # Sanitize the topic name to be suitable for folder names
    sanitized_topic = re.sub(r'[^\w\s-]', '', topic_name).strip()
    return sanitized_topic

def get_timestamp():
    """Generate a timestamp string for filenames."""
    return datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

def get_ai_response(prompt=None, user_id=None):
    """Get AI response for the main content."""
    global topic, sections, topic_directory, user_directory
    
    if user_id is None:
        # Fallback to global variables for backward compatibility
        current_topic = topic
        current_sections = sections
        output_dir = Path("outputFiles")
        if topic_directory:
            output_dir = topic_directory
        elif user_directory and current_topic:
            sanitized_topic = re.sub(r'[^\w\s-]', '', current_topic).strip().replace(' ', '_')
            output_dir = Path(user_directory) / sanitized_topic
    else:
        # Use user session data
        user_session = get_user_session(user_id)
        current_topic = user_session.topic
        current_sections = user_session.sections
        user_directory = user_session.user_directory
        
        # Determine output directory using user session
        output_dir = Path("outputFiles")
        if user_directory and current_topic:
            sanitized_topic = re.sub(r'[^\w\s-]', '', current_topic).strip().replace(' ', '_')
            output_dir = Path(user_directory) / sanitized_topic
            # Ensure directory exists
            output_dir.mkdir(parents=True, exist_ok=True)
    
    # If we already have the sections, return them
   
        
    if not current_topic:
        raise ValueError("Topic cannot be empty")
        
    try:
        # Store response in sections variable
        response = ai_agent(prompt)
        
        # Update the appropriate sections storage
        if user_id is not None:
            user_session.sections = response
            current_sections = response
        else:
            # Update global sections for backward compatibility
            sections = response
            current_sections = response
        
        # Ensure directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Add timestamp to filename
        timestamp = get_timestamp()
        output_file = output_dir / f"Slides_{timestamp}.txt"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(current_sections)
        
        print(f"AI response saved to {output_file}")
        return current_sections
    
    except Exception as e:
        print(f"Error in get_AI_response: {str(e)}")
        raise

def create_questions(custom_prompt, output_formats, user_id):
    """Generate questions document and webpage.
    
    Args:
        custom_prompt (str, optional): Custom prompt to generate content.
        output_formats (dict, optional): Dictionary with keys 'docx', 'html' and boolean values.
                                        Defaults to {'docx': True, 'html': False}.
        user_id: The ID of the current user.
    
    Returns:
        dict or str: Dictionary with paths to generated files or the path to the docx file for backward compatibility.
    """
    if user_id is None:
        raise ValueError("user_id is required")
        
    # Get user-specific session
    user_session = get_user_session(user_id)
    topic = user_session.topic
    learning_outcomes = user_session.learning_outcomes
    user_directory = user_session.user_directory
    
    print(f"Topic directory: {topic}")
    print(f"Generating Questions about ({topic})...")
    
    # Set default output formats if not specified
    if output_formats is None:
        output_formats = {'docx': True, 'html': False}
    
    if custom_prompt:
        prompt = custom_prompt
    else:
        try:
            # Try to read from questions.txt as fallback
            prompt = read_file_contents("questions.txt") + f" about: {topic}"
            if learning_outcomes:
                prompt += f" considering these learning outcomes: {learning_outcomes}"
        except Exception as e:
            print(f"Warning: Could not read from questions.txt: {str(e)}")
            # Fallback to a basic prompt if file can't be read
            prompt = f"Create comprehensive questions about {topic}"
            if learning_outcomes:
                prompt += f" considering these learning outcomes: {learning_outcomes}"
                
    questions = ai_agent(prompt)
    
    # Determine output directory
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic_directory if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
        # Update topic_directory global variable
        topic_directory = output_dir
    elif topic_directory:
        output_dir = topic_directory
    
    # Dictionary to store results
    result = {"content": questions}
    
    # Generate timestamp for filenames
    timestamp = get_timestamp()
    
    # Create document if DOCX format is selected
    if output_formats.get('docx', True):
        doc = Document()
        doc.add_heading('Questions', level=1)
        doc.add_paragraph(questions)
        
        # Save to the topic directory with timestamp
        output_file = output_dir / f"Questions_{timestamp}.docx"
        doc.save(str(output_file))
        print(f"Questions document saved to {output_file}")
        result["docx_path"] = str(output_file)

    # Generate HTML if HTML format is selected
    if output_formats.get('html', False):
        print("Generating questions webpage...")
        html_content = ai_agent(DESIGN + "<data>" + questions + "</data>. use the data as it is don't change any thing.")
        html_content = html_content.replace('```html', '').replace('```', '')
        
        # Save HTML to the topic directory with timestamp
        output_html = output_dir / f"questions_{timestamp}.html" 
        output_html.write_text(html_content, encoding="utf-8")
        print(f"Questions webpage saved to {output_html}")
        result["html_path"] = str(output_html)
    
    # For backward compatibility, if only docx is generated, return just the path
    if 'docx_path' in result and len(result) == 2:  # content + docx_path only
        return result["docx_path"]
    
    return result


def create_worksheet(custom_prompt=None, output_formats=None, user_id=None):
    """Generate worksheet document and webpage based on specified output formats."""
    if user_id is None:
        raise ValueError("user_id is required")
        
    # Get user-specific session
    user_session = get_user_session(user_id)
    topic = user_session.topic
    learning_outcomes = user_session.learning_outcomes
    user_directory = user_session.user_directory
    
    print(f"Generating worksheet about ({topic})...")
    
    # Default to DOCX only if no formats specified
    if output_formats is None:
        output_formats = {'docx': True, 'html': False}
    
    if custom_prompt:
        prompt = custom_prompt
    else:
        try:
            prompt = read_file_contents("worksheet.txt") + f" about: {topic}"
            if learning_outcomes:
                prompt += f" considering these learning outcomes: {learning_outcomes}"
            prompt += " give me the worksheet directly without explination at the begining of your response."
        except Exception as e:
            print(f"Warning: Could not read from worksheet.txt: {str(e)}")
            # Fallback to a basic prompt if file can't be read
            prompt = f"Create a comprehensive worksheet about {topic}"
            if learning_outcomes:
                prompt += f" considering these learning outcomes: {learning_outcomes}"
    
    worksheet = ai_agent(prompt)
    worksheet = worksheet.replace("**", "")
    
    # Determine output directory using user session
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
    
    # Dictionary to store results
    result = {"content": worksheet}
    
    # Generate ONE timestamp for ALL filenames to avoid duplicates
    timestamp = get_timestamp()
    
    # Create document if DOCX format is selected
    if output_formats.get('docx', True):
        doc = Document()
        doc.add_heading('Worksheet', level=1)
        doc.add_paragraph(worksheet)
        
        # Save to the topic directory with timestamp
        output_file = output_dir / f"worksheet_{timestamp}.docx"
        doc.save(str(output_file))
        print(f"Worksheet document saved to {output_file}")
        result["docx_path"] = str(output_file)
        result["file_path"] = str(output_file)  # For backward compatibility

    # Generate HTML if HTML format is selected
    if output_formats.get('html', False):
        print("Generating worksheet webpage...")
        html_content = ai_agent(DESIGN + "<data>" + worksheet + "</data>. use the data as it is don't change any thing.")
        html_content = html_content.replace('```html', '').replace('```', '')
        
        # Save HTML to the topic directory with timestamp
        output_html = output_dir / f"worksheet_{timestamp}.html"
        output_html.write_text(html_content, encoding="utf-8")
        print(f"Worksheet webpage saved to {output_html}")
        result["html_path"] = str(output_html)
    
    # For backward compatibility
    if not result.get("file_path") and result.get("html_path"):
        result["file_path"] = result["html_path"]
        
    return result


def create_task():
    """Generate task document and webpage."""
    global topic, topic_directory, user_directory
    print(f"Generating Task about ({topic})...")
    prompt = read_file_contents("task.txt") + f" about: {topic}"
    global learning_outcomes
    if learning_outcomes:
        prompt += f" considering these learning outcomes: {learning_outcomes}"
    prompt += " give me the task directly without explination at the begining of your response."
    
    task = ai_agent(prompt)
    task = task.replace("**", "")
    
    # Determine output directory
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic_directory if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
        # Update topic_directory global variable
        topic_directory = output_dir
    elif topic_directory:
        output_dir = topic_directory
    
    # Ensure directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate timestamp for filenames
    timestamp = get_timestamp()
    
    # Create document
    doc = Document()
    doc.add_heading('Task', level=1)
    doc.add_paragraph(task)
    
    # Save to the topic directory with timestamp
    output_file = output_dir / f"task_{timestamp}.docx"
    doc.save(str(output_file))
    print(f"Task document saved to {output_file}")

    print("Generating task webpage ...")
    html_content = ai_agent(DESIGN + "<data>" + task + "</data>. use the data as it is don't change any thing.")
    html_content = html_content.replace('```html', '').replace('```', '')
    
    # Save HTML to the topic directory with timestamp
    output_html = output_dir / f"task_{timestamp}.html"
    output_html.write_text(html_content, encoding="utf-8")
    print(f"Task webpage saved to {output_html}")


def create_homework(custom_prompt=None):
    """Generate homework document."""
    global topic, topic_directory, user_directory
    print(f"Generating homework about ({topic})...")
    
    if custom_prompt:
        prompt = custom_prompt
    else:
        prompt = read_file_contents("homework.txt") + f" about: {topic}"
        global learning_outcomes
        if learning_outcomes:
            prompt += f" considering these learning outcomes: {learning_outcomes}"
        prompt += " give me the homework directly without explination at the begining of your response."
    
    hw = ai_agent(prompt)
    hw = hw.replace("**", "")
    
    # Determine output directory
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic_directory if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
        # Update topic_directory global variable
        topic_directory = output_dir
    elif topic_directory:
        output_dir = topic_directory
    
    # Ensure directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Generate timestamp for filenames
    timestamp = get_timestamp()
    
    # Create document
    doc = Document()
    doc.add_heading('Homework', level=1)
    doc.add_paragraph(hw)
    
    # Save to the topic directory with timestamp
    output_file = output_dir / f"homework_{timestamp}.docx"
    doc.save(str(output_file))
    print(f"Homework saved to {output_file}")


def create_revision(custom_prompt=None, output_formats=None, user_id=None):
    """Generate revision document and webpage.
    
    Args:
        custom_prompt (str, optional): Custom prompt to generate content.
        output_formats (dict, optional): Dictionary with keys 'docx', 'html' and boolean values.
                                        Defaults to {'docx': True, 'html': True}.
        user_id: The ID of the current user.
    
    Returns:
        dict: Dictionary with paths to generated files
    """
    if user_id is None:
        raise ValueError("user_id is required")
        
    # Get user-specific session
    user_session = get_user_session(user_id)
    topic = user_session.topic
    learning_outcomes = user_session.learning_outcomes
    user_directory = user_session.user_directory
    
    print(f"Generating revision about ({topic})...")
    
    # Set default output formats if not specified
    if output_formats is None:
        output_formats = {'docx': True, 'html': True}
    
    if custom_prompt:
        prompt = custom_prompt
    else:
        prompt = read_file_contents("revision.txt")
        if learning_outcomes:
            prompt += f" considering these learning outcomes: {learning_outcomes}"
    
    hw = ai_agent(prompt)
    hw = hw.replace("**", "")
    
    # Determine output directory
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic_directory if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
        # Update topic_directory global variable
        global topic_directory
        topic_directory = output_dir
    elif topic_directory:
        output_dir = topic_directory
    
    # Ensure directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Dictionary to store results
    result = {"content": hw}
    
    # Generate timestamp for filenames
    timestamp = get_timestamp()
    
    # Create document if DOCX format is selected
    if output_formats.get('docx', True):
        doc = Document()
        doc.add_heading('Revision', level=1)
        doc.add_paragraph(hw)
        
        # Save to the topic directory with timestamp
        output_file = output_dir / f"revision_{timestamp}.docx"
        doc.save(str(output_file))
        print(f"Revision document saved to {output_file}")
        result["docx_path"] = str(output_file)

    # Generate HTML if HTML format is selected
    if output_formats.get('html', True):
        print("Generating revision webpage ...")
        html_content = ai_agent(DESIGN + "<data>" + hw + "</data>. use the data as it is don't change any thing.")
        html_content = html_content.replace('```html', '').replace('```', '')
        
        # Save HTML to the topic directory with timestamp
        output_html = output_dir / f"revision_{timestamp}.html" 
        output_html.write_text(html_content, encoding="utf-8")
        print(f"Revision webpage saved to {output_html}")
        result["html_path"] = str(output_html)
    
    return result


def create_lesson_plan(custom_prompt=None, output_formats=None, user_id=None):
    """Generate a comprehensive lesson plan based on the topic and learning outcomes.
    This function handles both the document creation and HTML generation."""
    
    if user_id is None:
        raise ValueError("user_id is required")
        
    # Get user-specific session
    user_session = get_user_session(user_id)
    topic = user_session.topic
    learning_outcomes = user_session.learning_outcomes
    user_directory = user_session.user_directory
    
    # Set default output formats if not specified
    if output_formats is None:
        output_formats = {'docx': True, 'pdf': True, 'html': True}
    
    print(f"Generating comprehensive lesson plan for ({topic})...")
    
    if custom_prompt:
        prompt = custom_prompt
    else:
        prompt = f"Create a detailed lesson plan for the topic: {topic}."
        if learning_outcomes:
            prompt += f"\n\nEnsure the lesson plan addresses these learning outcomes:\n{learning_outcomes}"
        
        prompt += "\n\nInclude the following sections in the lesson plan:\n"
        prompt += "- Learning Objectives / Outcomes\n"
        prompt += "- Materials / Resources Needed\n"
        prompt += "- Introduction / Warm-Up\n"
        prompt += "- Instructional Steps / Procedure\n"
        prompt += "- Guided Practice\n"
        prompt += "- Independent Practice\n"
        prompt += "- Assessment / Evaluation\n"
        prompt += "- Closure\n"
    
    # Generate the lesson plan content using the AI agent
    lesson_plan_content = ai_agent(prompt)
    # Remove asterisks from the AI-generated content
    clean_lesson_plan_content = lesson_plan_content.replace('*', '')
    
    # Determine output directory
    output_dir = Path("outputFiles")
    
    # Use the user_directory and topic_directory if they exist
    if user_directory and topic:
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        output_dir = Path(user_directory) / sanitized_topic
        # Ensure directory exists
        output_dir.mkdir(parents=True, exist_ok=True)
        # Update topic_directory global variable
        global topic_directory
        topic_directory = output_dir
    elif topic_directory:
        output_dir = topic_directory
    
    # Ensure directory exists
    os.makedirs(output_dir, exist_ok=True)
    
    # Save the clean lesson plan content to a text file
    timestamp = get_timestamp()
    lessonplan_txt_file = output_dir / f"lessonplan_{timestamp}.txt"
    with open(lessonplan_txt_file, 'w', encoding='utf-8') as f:
        f.write(clean_lesson_plan_content)
    print(f"Lesson plan content saved to {lessonplan_txt_file}")
    
    # Create document for the lesson plan if DOCX is selected
    result = {"content": lesson_plan_content}
    
    if output_formats.get('docx', True):
        timestamp = get_timestamp()
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        
        # Check if a template is specified
        template_name = output_formats.get('docx_template', None)
        
        if template_name == 'template1':
            # Use the template1.docx
            template_path = Path("word templates") / "template1.docx"
            
            if template_path.exists():
                print(f"Using template: {template_path}")
                
                # Extract individual sections from the content
                section_pattern = r'(?:^|\n)(#+\s*([^#\n]+))\s*\n((?:.+\n?)+?)(?=\n#+\s+|$)'
                sections = re.findall(section_pattern, clean_lesson_plan_content, re.MULTILINE)
                
                # Convert to dictionary of sections
                sections_dict = {}
                
                # First, check for Learning Objectives
                learning_obj_pattern = r'(?:Learning Objectives|Learning Outcomes)[^\n]*\n((?:.+\n?)+?)(?=\n#+\s+|$)'
                learning_obj_match = re.search(learning_obj_pattern, clean_lesson_plan_content, re.IGNORECASE | re.MULTILINE)
                
                if learning_obj_match:
                    sections_dict["Learning Objectives"] = learning_obj_match.group(1).strip()
                
                # Print Learning Objectives section if found
                
                
                # Extract individual sections from the content using string splitting
                sections_data = clean_lesson_plan_content.split("!!!")
                
                # Convert to dictionary of sections
                sections_dict = {}
                
                # Add global learning outcomes to the dictionary if available
                if learning_outcomes:
                    sections_dict["Learning Objectives"] = learning_outcomes
                # Add global topic to the dictionary if available
                if topic:
                    sections_dict["topic"] = topic
                # Process each section
                for section in sections_data:
                    if not section.strip():
                        continue
                        
                    # Try to split into title and content
                    lines = section.strip().split('\n', 1)
                    if len(lines) >= 2:
                        section_title = lines[0].strip()
                        section_content = lines[1].strip()
                        
                        # Print the current section being processed
                       # print(f"\n=== SECTION: {section_title} ===")
                        #print(section_content[:200] + "..." if len(section_content) > 200 : section_content)
                        
                        # Map section titles to template section names
                        if re.search(r'Learning Objectives|Learning Outcomes', section_title, re.IGNORECASE):
                            sections_dict["Learning Objectives"] = section_content
                        elif re.search(r'material|resource', section_title, re.IGNORECASE):
                            sections_dict["Resources Needed"] = section_content
                            print("Resources Needed section found and added.")
                        elif re.search(r'introduction|warm[- ]?up', section_title, re.IGNORECASE):
                            sections_dict["Introduction"] = section_content
                        elif re.search(r'instructional|procedure|steps', section_title, re.IGNORECASE):
                            sections_dict["Instructional Steps"] = section_content
                        elif re.search(r'guided\s+practice', section_title, re.IGNORECASE):
                            sections_dict["Guided Practice"] = section_content
                        elif re.search(r'independent\s+practice', section_title, re.IGNORECASE):
                            sections_dict["Independent Practice"] = section_content
                        elif re.search(r'Assessment / Evaluation', section_title, re.IGNORECASE):
                            sections_dict["Assessment / Evaluation"] = section_content
                        elif re.search(r'closure', section_title, re.IGNORECASE):
                            sections_dict["Closure"] = section_content
                        elif re.search(r'differentiation', section_title, re.IGNORECASE):
                            sections_dict["Differentiation"] = section_content
                        elif re.search(r'Teacher Directed Activities', section_title, re.IGNORECASE):
                            sections_dict["Teacher Directed Activities"] = section_content
                        elif re.search(r'Student Directed Activities', section_title, re.IGNORECASE):
                            sections_dict["Student Directed Activities"] = section_content
                        elif re.search(r'Extension, Refinement, and Practice Activities', section_title, re.IGNORECASE):
                            sections_dict["Extension, Refinement, and Practice Activities"] = section_content
                        elif re.search(r'Use of Media & Technology', section_title, re.IGNORECASE):
                            sections_dict["Use of Media / Technology"] = section_content
                        elif re.search(r'Cooperative Grouping', section_title, re.IGNORECASE):
                            sections_dict["Cooperative Grouping"] = section_content
                        elif re.search(r'Differentiation', section_title, re.IGNORECASE):
                            sections_dict["Differentiation"] = section_content
                        else:
                            # Generic section handling   Use of Media & Technology
                            sections_dict[section_title] = section_content
                
                # Print the final mapping of sections
                #print("\n=== FINAL SECTIONS DICTIONARY ===")
                #for section_name, content in sections_dict.items():
                 #   print(f"Section: {section_name}")
                  #  print(f"Content preview: {content[:100]}..." if len(content) > 100 : content)
                   # print("-" * 50)
                
                # Populate the template
                doc = populate_template_with_sections(template_path, sections_dict, topic)
                
                if not doc:
                    # Fallback to standard document creation if template population fails
                    doc = Document()
                    doc.add_heading(f'Lesson Plan: {topic}', level=1)
                    doc.add_paragraph(clean_lesson_plan_content)
            else:
                print(f"Template not found at {template_path}. Using standard document instead.")
                doc = Document()
                doc.add_heading(f'Lesson Plan: {topic}', level=1)
                doc.add_paragraph(clean_lesson_plan_content)
        else:
            # Standard document creation (no template)
            doc = Document()
            doc.add_heading(f'Lesson Plan: {topic}', level=1)
            doc.add_paragraph(clean_lesson_plan_content)
        
        # Save to the topic directory with timestamp
        output_file = output_dir / f"LessonPlan_{timestamp}.docx"
        doc.save(str(output_file))
        print(f"Comprehensive lesson plan saved to {output_file}")
        result["docx_path"] = str(output_file)
    
    # Format the same lesson plan content as HTML if HTML is selected
    if output_formats.get('html', True):
        # We don't generate new content with AI here, just use the existing content
        html_content = ai_agent(DESIGN + "<data>" + clean_lesson_plan_content + "</data>. use the data as it is don't change any thing.")
        html_content = html_content.replace('```html', '').replace('```', '')
        
        sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
        html_filename = f"LessonPlan_{timestamp}.html"
        output_html = output_dir / html_filename
        
        # Save HTML to the topic directory
        output_html.write_text(html_content, encoding="utf-8")
        print(f"Lesson plan webpage saved to {output_html}")
        result["html_path"] = str(output_html)
    
    return result

def populate_template_with_sections(template_path, sections_content, topic_title):
    """
    Populate a Word template with lesson plan sections.
    
    Args:
        template_path (str): Path to the template DOCX file
        sections_content (dict): Dictionary of section names and their content
        topic_title (str): The title/topic of the lesson plan
        
    Returns:
        Document: Populated Word document
    """
    try:
        # Load the template
        doc = Document(template_path)
        
        # Find the table in the template
        target_table = None
        for table in doc.tables:
            target_table = table
            break
            
        if not target_table:
            print("Warning: No table found in template. Creating standard document instead.")
            return None
            
        # Set the title in the document
        for paragraph in doc.paragraphs:
            if "TITLE_PLACEHOLDER" in paragraph.text:
                paragraph.text = paragraph.text.replace("TITLE_PLACEHOLDER", topic_title)
                break
        
        # Iterate through all cells in the table to find specific placeholders
        print("Processing table cells for placeholders...")
        for row in target_table.rows:
            for cell in row.cells:
                cell_text = cell.text
                #print(f"Processing cell: {cell_text}")
                # Check for Learning Outcomes placeholder (LO111)
                if "LO111" in cell_text and "Learning Objectives" in sections_content:
                    # Get the learning objectives content
                    cleaned_objectives = sections_content["Learning Objectives"]
                    
                    # First, normalize line breaks by replacing all types of line breaks with \n
                    cleaned_objectives = cleaned_objectives.replace('\r\n', '\n').replace('\r', '\n')
                    
                    # Split by lines, remove empty lines, and trim whitespace
                    lines = [line.strip() for line in cleaned_objectives.split('\n') if line.strip()]
                    
                    # Rejoin with single line breaks
                    cleaned_objectives = '\n'.join(lines)
                    
                    # Replace the placeholder with cleaned content
                    cell.text = cell_text.replace("LO111", cleaned_objectives)
                    print("Replaced LO111 with Learning Objectives content.")
                
                # Check for Warm Up / Introduction placeholder (WU111) 
                elif "WU111" in cell_text and "Introduction" in sections_content:
                    cell.text = cell_text.replace("WU111", sections_content["Introduction"])
                    print("Replaced WU111 with Introduction content.")
                    
                
                 # Check for Topic placeholder (TO111)
                elif "TO111" in cell_text and "topic" in sections_content:
                    cell.text = cell_text.replace("TO111", sections_content["topic"])
                    print("Replaced TO111 with topic content.")
                    
                    # Check for Date placeholder (DA111)
                elif "DA111" in cell_text:
                    cell.text = cell_text.replace("DA111", datetime.datetime.now().strftime("%B %d, %Y"))
                    print("Replaced DA111 with current date.")

                 # Check for Teacher Directed Activities placeholder (TO111)
                elif "TDAC111" in cell_text and "Teacher Directed Activities" in sections_content:
                    cell.text = cell_text.replace("TDAC111", sections_content["Teacher Directed Activities"])
                    print("Replaced TDAC111 with Teacher Directed Activities content.")
                # Check for Student Directed Activities placeholder (SDA111)
                elif "SDAC111" in cell_text and "Student Directed Activities" in sections_content:
                    cell.text = cell_text.replace("SDAC111", sections_content["Student Directed Activities"])
                    print("Replaced SDAC111 with Student Directed Activities content.")
                
                # Check for Extension, Refinement, and Practice Activities placeholder (ERPA111)
                elif "ERPA111" in cell_text and "Extension, Refinement, and Practice Activities" in sections_content:
                    cell.text = cell_text.replace("ERPA111", sections_content["Extension, Refinement, and Practice Activities"])
                    print("Replaced ERPA111 with Extension, Refinement, and Practice Activities content.")
                
                # check for Assessment / Evaluation placeholder (FA111)
                elif "FA111" in cell_text and "Assessment / Evaluation" in sections_content:
                    cell.text = cell_text.replace("FA111", sections_content["Assessment / Evaluation"])
                    print("Replaced FA111 with Assessment / Evaluation content.")

                # check for Assessment / Evaluation placeholder (FA111)
                elif "CA111" in cell_text and "Closure" in sections_content:
                    cell.text = cell_text.replace("CA111", sections_content["Closure"])
                    print("Replaced CA111 with Closure content.")

                # Check for Resources Needed placeholder (RN111)
                
                elif "RN111" in cell_text and "Resources Needed" in sections_content:
                    cell.text = cell_text.replace("RN111", sections_content["Resources Needed"])
                    print("Replaced RN111 with Resources Needed content.")
                    

                # Check for Use of Media / Technology placeholder (UOMT111)
                elif "UOMT111" in cell_text and "Use of Media / Technology" in sections_content:
                    cell.text = cell_text.replace("UOMT111", sections_content["Use of Media / Technology"])
                    print("Replaced UOMT111 with Use of Media / Technology content.")
                
                # check for Cooperative Grouping placeholder (CG111)
                elif "CG111" in cell_text and "Cooperative Grouping" in sections_content:
                    cell.text = cell_text.replace("CG111", sections_content["Cooperative Grouping"])
                    print("Replaced CG111 with Cooperative Grouping content.")
                
                # check for Differentiation placeholder (DIFF111)
                elif "DIFF111" in cell_text and "Differentiation" in sections_content:
                    cell.text = cell_text.replace("DIFF111", sections_content["Differentiation"])
                    print("Replaced DIFF111 with Differentiation content.")
                
        # Continue with matching section names to table cells for other sections
        for i, row in enumerate(target_table.rows):
            if i == 0:  # Skip header row if it exists
                continue
                
            if len(row.cells) < 2:
                continue
                
            section_cell = row.cells[0]
            content_cell = row.cells[1]
            
            section_name = section_cell.text.strip()
            
            # Skip cells that have already been processed with our placeholders
            if "LO111" in section_name or "WU111" in section_name:
                continue
                
            # Check if this section has content in our dictionary
            if section_name in sections_content:
                content_cell.text = sections_content[section_name]
                
        # Add date to the document
        current_date = datetime.datetime.now().strftime("%B %d, %Y")
        for paragraph in doc.paragraphs:
            if "DATE_PLACEHOLDER" in paragraph.text:
                paragraph.text = paragraph.text.replace("DATE_PLACEHOLDER", current_date)
                
        return doc
        
    except Exception as e:
        print(f"Error populating template: {str(e)}")
        traceback.print_exc()
        return None

def apply_theme_to_slide(slide, theme_name):
    """Apply theme-specific styling to a slide."""
    # Skip applying any theme to maintain plain formatting
    print(f"Skipping theme application due to no-formatting requirement")
    return True

def create_presentation(custom_prompt, theme_name, user_id=None):
    """Create PowerPoint presentation based on sections data.
    
    Args:
        custom_prompt (str, optional): Custom prompt to generate content if sections is None.
        theme_name (str, optional): Name of the PowerPoint theme to apply. Defaults to "Office".
        user_id: The ID of the current user.
        
    Returns:
        str: Path to the saved PowerPoint file.
    """
    if user_id is None:
        raise ValueError("user_id is required")
        
    # Get user-specific session
    user_session = get_user_session(user_id)
    topic = user_session.topic
    #sections = user_session.sections
    learning_outcomes = user_session.learning_outcomes
    user_directory = user_session.user_directory
    
    # Get AI response for content if not already available
    
    get_ai_response(custom_prompt, user_id=user_id)
    # Update sections in user session (get_ai_response should have updated it)
    sections = user_session.sections
    
    # Remove all asterisks from sections
    sections = sections.replace('*', '') if sections else None
    
    print(f"Creating PowerPoint presentation for topic: {topic} with theme: {theme_name}")
    try:
        # Check for theme-specific template in the "powerpoint templates" folder first
        theme_template_path = Path("powerpoint templates") / f"{theme_name}.pptx"
        
        
        presentation = None
        
        # Try to use the theme-specific template first
        if theme_template_path.exists():
            print(f"Using theme-specific template: {theme_template_path}")
            try:
                presentation = Presentation(theme_template_path)
                print(f"Successfully loaded theme template: {theme_template_path}")
            except Exception as e:
                print(f"Error loading theme template: {e}")
        
        
        
        # If both templates failed, create a blank presentation
        if presentation is None:
            print("No templates available, creating blank presentation")
            presentation = Presentation()
            
        # Create the title slide (first slide)
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        
        # Attempt to get title and subtitle placeholders
        try:
            title_shape = title_slide.shapes.title
            subtitle_shape = title_slide.placeholders[1]
        except (KeyError, IndexError, AttributeError) as e:
            print(f"Error getting title slide placeholders: {e}")
            # Create manual text boxes if placeholders not available
            top = Inches(1)
            left = Inches(0.5)
            width = Inches(9)
            height = Pt(100)
            title_shape = title_slide.shapes.add_textbox(left, top, width, height)
        
        # Set title - limit length to avoid corruption
        safe_topic = topic[:100] if topic and len(topic) > 100 else topic
        if hasattr(title_shape, 'text_frame'):
            title_shape.text_frame.text = safe_topic
            # No font formatting applied
        
        # Add date to subtitle if placeholder exists
        if subtitle_shape and hasattr(subtitle_shape, 'text_frame'):
            today = datetime.datetime.now().strftime("%B %d, %Y")
            subtitle_shape.text_frame.text = f"Created: {today}"
            # No font formatting applied
        print(f"Sections____: {sections}____________")
        # Process the sections from AI-generated content
        if sections:
            sections_list = [part.strip() for part in sections.split('!!!') if part.strip()]
            
            # Extract titles using the utility function
            titles_list = extract_important_lines(sections)
            titles_list = [title.replace("!!!", "").strip() for title in titles_list]
            print(f"Extracted {len(titles_list)} titles for slides.")
            # Create slides for each section
            for i, (section_content, title) in enumerate(zip(sections_list, titles_list)):
                print(f"Processing slide {i+1} with title: {title}")
                if not title or not section_content:
                    continue
                
                # Create a new slide
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                title_shape = slide.shapes.title
                
                # Clean title and remove any meta formatting
                clean_title = title
                # Remove any "Section X: " or "Information X: " prefixes
                clean_title = re.sub(r'^(Section|Information|Slide)\s+\d+:\s*', '', clean_title)
                # Remove any "(Slide X)" suffix
                clean_title = re.sub(r'\s*\(Slide\s+\d+\)\s*$', '', clean_title)
                # Remove any "**Slide Title:**" prefixes and other markdown formatting
                clean_title = re.sub(r'\*\*Slide\s+Title:\*\*\s*', '', clean_title)
                clean_title = clean_title.replace('*', '')
                
                # Look for customized slide title in the content
                custom_title_match = re.search(r'Title:\s*([^"\n]+)', section_content)
                if custom_title_match:
                    custom_title = custom_title_match.group(1).strip()
                    if custom_title:
                        clean_title = custom_title
                
                # Limit title length to avoid corruption
                safe_title = clean_title[:100] if clean_title and len(clean_title) > 100 else clean_title
                title_shape.text = safe_title
                
                # No font formatting applied to title
                
                # Process content for the slide
                if len(slide.placeholders) > 1:
                    content_shape = slide.placeholders[1]
                    content_text_frame = content_shape.text_frame
                    
                    # Clean and prepare content
                    clean_content = section_content.replace('\r', '').strip()
                    
                    # Remove the title from the content
                    clean_content = clean_content.replace(title, '', 1).strip()
                    
                    # Remove "Title:..." text from content
                    clean_content = re.sub(r'Title:\s*([^"\n]+)', '', clean_content).strip()
                    
                    # Look for image placeholder text that should go to notes
                    image_placeholder_pattern = r'(A captivating image of[^.]*\.)'
                    image_placeholders = re.findall(image_placeholder_pattern, clean_content)
                    
                    # Remove image placeholders from content
                    for placeholder in image_placeholders:
                        clean_content = clean_content.replace(placeholder, '').strip()
                    
                    # Remove content markers like "**Content:**", "**Speaker Notes:**", etc.
                    clean_content = re.sub(r'\*\*Content:\*\*', '', clean_content)
                    clean_content = re.sub(r'\*\*Speaker\s+Notes:\*\*', '', clean_content)
                    clean_content = re.sub(r'\*\*Visual:\*\*', '', clean_content)
                    
                    # Remove any markdown formatting
                    clean_content = clean_content.replace('**', '')
                    
                    # Split the content by lines to handle line breaks properly
                    content_length = len(clean_content)
                    if content_length > 1000:  # Limit content length
                        clean_content = clean_content[:1000] + "..."
                    
                    if content_length < 200:  # Short content
                        font_size = 24
                    elif content_length < 500:  # Medium content
                        font_size = 18
                    elif content_length < 800:  # Long content
                        font_size = 14
                    else:  # Very long content
                        font_size = 12
                    # Add the content paragraph by paragraph
                    paragraphs = clean_content.split('\n')
                    for j, paragraph_text in enumerate(paragraphs):
                        if not paragraph_text.strip():
                            continue
                        
                        p = content_text_frame.add_paragraph() if j > 0 else content_text_frame.paragraphs[0]
                        p.text = paragraph_text.strip()
                        p.font.size = Pt(font_size)
                        p.font.name = 'Arial'
                        p.font.bold = True
                        # No font formatting applied
                    
                    # Add image placeholders and other notes to the notes section
                    if hasattr(slide, 'notes_slide') and image_placeholders:
                        notes_slide = slide.notes_slide
                        notes_text_frame = notes_slide.notes_text_frame
                        
                        for placeholder in image_placeholders:
                            notes_text_frame.text += placeholder + "\n"
        
        # Do not apply theme-specific styling to slides
        # Skip the theme application entirely
        
        # Determine output directory - Fixed to use proper user directory structure
        output_dir = Path("outputFiles")
        
        # Use the user_directory and topic_directory if they exist
        if user_directory and topic:
            sanitized_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
            output_dir = Path(user_directory) / sanitized_topic
            # Ensure directory exists
            output_dir.mkdir(parents=True, exist_ok=True)
            print(f"Using user-specific directory: {output_dir}")
        else:
            print(f"Warning: Using default outputFiles directory. user_directory={user_directory}, topic={topic}")
        
        # Ensure directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Create sanitized file name based on topic - keep it simple
        safe_topic_name = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')        # Limit filename length to avoid path length issues
        if len(safe_topic_name) > 20:
            safe_topic_name = safe_topic_name[:20]
            
        if not safe_topic_name:
            safe_topic_name = "Lesson_Presentation"
            
        # Add timestamp to ensure uniqueness
        timestamp = get_timestamp()
        ppt_filename = f"Presentation_{timestamp}.pptx"
        ppt_path = output_dir / ppt_filename
        
        # Save the file and ensure it was created
        presentation.save(str(ppt_path))
        if not ppt_path.exists():
            print(f"Failed to create PowerPoint file at: {ppt_path}")
            return None
            
        print(f"Success! Created PowerPoint file: {ppt_path}")
        return str(ppt_path)
        
    except Exception as e:
        print(f"Error creating PowerPoint presentation: {e}")
        traceback.print_exc()
        return None

def extract_content(file_path_or_stream):
    """
    Extract text content from various file formats (.txt, .docx, .pdf, .pptx).
    
    Args:
        file_path_or_stream: Either a file path string or a file stream object
    
    Returns:
        str: Extracted text content from the file
    
    Raises:
        ValueError: If file format is not supported
        Exception: If file cannot be read or processed
    """
    try:
        # Determine if input is a file path or stream
        if hasattr(file_path_or_stream, 'read'):
            # It's a file stream
            file_stream = file_path_or_stream
            # Try to get filename from stream if available
            filename = getattr(file_stream, 'filename', getattr(file_stream, 'name', 'unknown'))
        else:
            # It's a file path
            file_path = Path(file_path_or_stream)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            filename = file_path.name
            file_stream = open(file_path, 'rb')
        
        # Determine file extension
        file_extension = Path(filename).suffix.lower()
        
        try:
            if file_extension == '.txt':
                # Handle text files
                if hasattr(file_path_or_stream, 'read'):
                    content = file_stream.read()
                    if isinstance(content, bytes):
                        content = content.decode('utf-8', errors='ignore')
                else:
                    with open(file_path_or_stream, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                return content.strip()
            
            elif file_extension == '.docx':
                # Handle Word documents
                if hasattr(file_path_or_stream, 'read'):
                    # Create a BytesIO object from the stream
                    file_bytes = file_stream.read()
                    doc_stream = io.BytesIO(file_bytes)
                    doc = Document(doc_stream)
                else:
                    doc = Document(file_path_or_stream)
                
                # Extract text from all paragraphs
                paragraphs = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text.strip())
                
                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                paragraphs.append(cell.text.strip())
                
                return '\n'.join(paragraphs)
            
            elif file_extension == '.pdf':
                # Handle PDF files
                if hasattr(file_path_or_stream, 'read'):
                    pdf_reader = PyPDF2.PdfReader(file_stream)
                else:
                    with open(file_path_or_stream, 'rb') as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                
                text_content = []
                for page in pdf_reader.pages:
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(page_text.strip())
                    except Exception as e:
                        print(f"Warning: Could not extract text from a PDF page: {e}")
                        continue
                
                return '\n'.join(text_content)
            
            elif file_extension == '.pptx':
                # Handle PowerPoint files
                if hasattr(file_path_or_stream, 'read'):
                    file_bytes = file_stream.read()
                    ppt_stream = io.BytesIO(file_bytes)
                    prs = Presentation(ppt_stream)
                else:
                    prs = Presentation(file_path_or_stream)
                
                text_content = []
                for slide in prs.slides:
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            text_content.append(shape.text.strip())
                        
                        # Extract text from tables in slides
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        text_content.append(cell.text.strip())
                    
                    # Extract notes from slide
                    if hasattr(slide, 'notes_slide'):
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            text_content.append(f"Notes: {notes_text.strip()}")
                
                return '\n'.join(text_content)
            
            else:
                raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: .txt, .docx, .pdf, .pptx")
        
        finally:
            # Close file stream if we opened it
            if not hasattr(file_path_or_stream, 'read'):
                file_stream.close()
    
    except Exception as e:
        print(f"Error extracting content from file: {str(e)}")
        raise Exception(f"Failed to extract content from file: {str(e)}")

